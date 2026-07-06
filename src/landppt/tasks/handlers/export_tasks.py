"""Worker handlers for file-based exports."""

from __future__ import annotations

import logging
import os
import shutil
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ...services.background_tasks import TaskStatus, get_task_manager
from ...services.pdf_to_pptx_converter import get_pdf_to_pptx_converter
from ...services.pyppeteer_pdf_converter import get_pdf_converter
from ...services.speech_script_repository import SpeechScriptRepository
from ...services.storage import get_artifact_service
from ...web.route_modules.export_support import _generate_pdf_with_pyppeteer, _prepare_html_for_file_based_export
from ...web.route_modules.narration_routes import _build_narration_audio_export_zip
from ...web.route_modules.support import ppt_service
from ..registry import task_handler

logger = logging.getLogger(__name__)


async def _save_artifact(task, local_path: str, artifact_type: str, filename: str, content_type: str) -> str:
    metadata = task.metadata if isinstance(task.metadata, dict) else {}
    artifact = await get_artifact_service().save_file(
        local_path=local_path,
        user_id=int(metadata["user_id"]),
        project_id=str(metadata.get("project_id")) if metadata.get("project_id") else None,
        task_id=task.task_id,
        artifact_type=artifact_type,
        filename=filename,
        content_type=content_type,
    )
    return artifact.id


async def _project_for_task(task):
    metadata = task.metadata if isinstance(task.metadata, dict) else {}
    return await ppt_service.project_manager.get_project(
        str(metadata["project_id"]),
        user_id=int(metadata["user_id"]),
    )


async def _set_progress(task_id: str, progress: float) -> None:
    await get_task_manager().update_task_status_async(task_id, TaskStatus.RUNNING, progress=progress)


@task_handler("pdf_generation")
async def export_pdf(task) -> dict:
    project = await _project_for_task(task)
    if not project or not project.slides_data:
        return {"success": False, "error": "Project not found or PPT not generated yet"}
    if not get_pdf_converter().is_available():
        return {"success": False, "error": "PDF generation service unavailable. Please ensure Playwright is installed."}

    temp_pdf_path = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    try:
        await _set_progress(task.task_id, 5)
        success = await _generate_pdf_with_pyppeteer(project, temp_pdf_path, individual=False)
        if not success or not os.path.exists(temp_pdf_path):
            return {"success": False, "error": "PDF generation failed"}
        artifact_id = await _save_artifact(
            task,
            temp_pdf_path,
            "pdf_export",
            f"{project.topic}_PPT.pdf",
            "application/pdf",
        )
        return {"success": True, "artifact_id": artifact_id, "project_topic": project.topic, "slide_count": len(project.slides_data)}
    finally:
        if os.path.exists(temp_pdf_path):
            os.unlink(temp_pdf_path)


@task_handler("pdf_to_pptx_conversion")
async def export_pptx(task) -> dict:
    project = await _project_for_task(task)
    if not project or not project.slides_data:
        return {"success": False, "error": "Project not found or PPT not generated yet"}

    pdf_converter = get_pdf_converter()
    converter = get_pdf_to_pptx_converter()
    metadata = task.metadata if isinstance(task.metadata, dict) else {}
    await converter.set_user_id_async(int(metadata["user_id"]))
    if not pdf_converter.is_available():
        return {"success": False, "error": "PDF generation service unavailable. Please ensure Playwright is installed."}
    if not converter.is_available():
        return {"success": False, "error": "PPTX conversion service unavailable. Please ensure Apryse SDK is installed and licensed."}

    temp_pdf_path = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False).name
    temp_pptx_path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
    try:
        await _set_progress(task.task_id, 5)
        pdf_success = await _generate_pdf_with_pyppeteer(project, temp_pdf_path, individual=False)
        if not pdf_success:
            return {"success": False, "error": "PDF generation failed"}

        await _set_progress(task.task_id, 50)
        success, result = await converter.convert_pdf_to_pptx_async(temp_pdf_path, temp_pptx_path)
        if not success:
            return {"success": False, "error": result}

        try:
            repo = SpeechScriptRepository()
            scripts_list = await repo.get_current_speech_scripts_by_project(str(metadata["project_id"]))
            speech_scripts = {script.slide_index: script.script_content for script in scripts_list}
            repo.close()
            if speech_scripts:
                prs = Presentation(temp_pptx_path)
                for index, slide in enumerate(prs.slides):
                    if index in speech_scripts:
                        slide.notes_slide.notes_text_frame.text = speech_scripts[index]
                prs.save(temp_pptx_path)
        except Exception as exc:
            logger.warning("Failed to add speech scripts to PPTX: %s", exc)

        await _set_progress(task.task_id, 90)
        artifact_id = await _save_artifact(
            task,
            temp_pptx_path,
            "pptx_export",
            f"{project.topic}_PPT.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        return {"success": True, "artifact_id": artifact_id, "project_topic": project.topic, "slide_count": len(project.slides_data)}
    finally:
        for path in (temp_pdf_path, temp_pptx_path):
            if os.path.exists(path):
                os.unlink(path)


@task_handler("html_to_pptx_screenshot")
async def export_html_screenshot_pptx(task) -> dict:
    metadata = task.metadata if isinstance(task.metadata, dict) else {}
    project = await _project_for_task(task)
    slides = metadata.get("slides") or []
    if not project or not slides:
        return {"success": False, "error": "Project not found or no slides provided"}
    pdf_converter = get_pdf_converter()
    if not pdf_converter.is_available():
        return {"success": False, "error": "Screenshot service unavailable. Please ensure Playwright is installed."}

    temp_dir = tempfile.mkdtemp()
    temp_pptx_path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
    try:
        speech_scripts = {}
        try:
            repo = SpeechScriptRepository()
            scripts_list = await repo.get_current_speech_scripts_by_project(str(metadata["project_id"]))
            speech_scripts = {script.slide_index: script.script_content for script in scripts_list}
            repo.close()
        except Exception as exc:
            logger.warning("Failed to load speech scripts: %s", exc)

        html_files = []
        export_base_url = str(metadata.get("export_base_url") or "")
        for index, slide in enumerate(slides):
            html_file = os.path.join(temp_dir, f"slide_{index}.html")
            prepared_html = _prepare_html_for_file_based_export(slide.get("html_content", ""), export_base_url)
            with open(html_file, "w", encoding="utf-8") as file_obj:
                file_obj.write(prepared_html)
            html_files.append(html_file)

        screenshot_paths = []
        for index, html_file in enumerate(html_files):
            screenshot_path = os.path.join(temp_dir, f"slide_{index}.png")
            success = await pdf_converter.screenshot_html(
                html_file,
                screenshot_path,
                width=1280,
                height=720,
                optimize_for_static=True,
                stability_checks=1,
                stability_interval=0.2,
            )
            if success:
                screenshot_paths.append(screenshot_path)
            await _set_progress(task.task_id, 25 + ((index + 1) / max(len(html_files), 1)) * 55)

        if not screenshot_paths:
            return {"success": False, "error": "No screenshots were generated"}

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        for index, screenshot_path in enumerate(screenshot_paths):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(screenshot_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            if index in speech_scripts:
                slide.notes_slide.notes_text_frame.text = speech_scripts[index]
        prs.save(temp_pptx_path)

        artifact_id = await _save_artifact(
            task,
            temp_pptx_path,
            "pptx_export",
            f"{project.topic}_PPT.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        return {"success": True, "artifact_id": artifact_id, "project_topic": project.topic, "slide_count": len(slides)}
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
        if os.path.exists(temp_pptx_path):
            os.unlink(temp_pptx_path)


@task_handler("narration_audio_export")
async def export_narration_audio(task) -> dict:
    metadata = task.metadata if isinstance(task.metadata, dict) else {}
    project = await _project_for_task(task)
    if not project:
        return {"success": False, "error": "Project not found"}

    from ...services.narration_service import NarrationService

    language = str(metadata.get("language") or "zh").strip().lower() or "zh"
    provider = str(metadata.get("provider") or "auto").strip().lower() or "auto"
    await _set_progress(task.task_id, 5)
    items = await NarrationService(user_id=int(metadata["user_id"])).generate_project_slide_audios(
        project_id=str(metadata["project_id"]),
        slide_indices=None,
        provider=provider,
        language=language,
        voice=metadata.get("voice"),
        rate=str(metadata.get("rate") or "+0%"),
        reference_audio_path=metadata.get("reference_audio_path"),
        reference_text=str(metadata.get("reference_text") or ""),
        voice_prompt=str(metadata.get("voice_prompt") or ""),
        force_regenerate=bool(metadata.get("force_regenerate")),
        uploads_dir="uploads",
    )
    if not items:
        return {"success": False, "error": "No narration audio available for export"}

    await _set_progress(task.task_id, 78)
    zip_path = _build_narration_audio_export_zip(
        project_topic=project.topic,
        slides_data=project.slides_data or [],
        language=language,
        items=items,
    )
    try:
        await _set_progress(task.task_id, 96)
        artifact_id = await _save_artifact(
            task,
            zip_path,
            "narration_audio_export",
            f"{project.topic}_narration_audio_{language}.zip",
            "application/zip",
        )
        return {"success": True, "artifact_id": artifact_id, "language": language, "provider": provider, "count": len(items)}
    finally:
        if os.path.exists(zip_path):
            os.unlink(zip_path)


@task_handler("narration_video_export")
async def export_narration_video(task) -> dict:
    metadata = task.metadata if isinstance(task.metadata, dict) else {}
    project = await _project_for_task(task)
    if not project:
        return {"success": False, "error": "Project not found"}

    from ...services.video_export_service import NarrationVideoExportService

    language = str(metadata.get("language") or "zh").strip().lower() or "zh"
    fps = 60 if int(metadata.get("fps") or 30) == 60 else 30
    await _set_progress(task.task_id, 5)
    result = await NarrationVideoExportService().export_project_video(
        project=project,
        language=language,
        fps=fps,
        width=1920,
        height=1080,
        embed_subtitles=bool(metadata.get("embed_subtitles", True)),
        subtitle_style=metadata.get("subtitle_style"),
        render_mode=str(metadata.get("render_mode") or "live"),
        uploads_dir="uploads",
    )
    if not isinstance(result, dict) or not result.get("success"):
        return result if isinstance(result, dict) else {"success": False, "error": "Narration video export failed"}

    video_path = result.get("video_path")
    if not video_path or not Path(video_path).exists():
        return {"success": False, "error": "Narration video file not found"}

    artifact_id = await _save_artifact(
        task,
        str(video_path),
        "narration_video_export",
        f"{project.topic}_narration_{language}.mp4",
        "video/mp4",
    )
    return {**result, "artifact_id": artifact_id, "video_path": None}
