"""
Global Master Template Service for managing reusable master templates
"""

import json
import logging
import time
import base64
from collections import Counter
from typing import Dict, Any, List, Optional
from io import BytesIO
from sqlalchemy.exc import IntegrityError

from ...ai import get_ai_provider, get_role_provider, AIMessage, MessageRole
from ...ai.base import TextContent, ImageContent, MessageContentType
from ...core.config import ai_config
from ...database.service import DatabaseService
from ...database.database import AsyncSessionLocal
from ..prompts.system_prompts import SystemPrompts
from ..prompts.template_prompts import TemplatePrompts

# Configure logger for this module
logger = logging.getLogger(__name__)


class GlobalMasterTemplateService:
    """Service for managing global master templates"""

    def __init__(
        self,
        provider_name: Optional[str] = None,
        user_id: Optional[int] = None,
        allow_system_template_write: bool = False,
    ):
        self.provider_name = provider_name
        self.user_id = user_id
        self.allow_system_template_write = bool(allow_system_template_write)

    @property
    def ai_provider(self):
        """Dynamically get AI provider to ensure latest config"""
        provider, _ = get_role_provider("template", provider_override=self.provider_name)
        return provider

    def _get_template_role_provider(self):
        """Get provider and settings for template generation role (sync version)"""
        return get_role_provider("template", provider_override=self.provider_name)

    async def _get_template_role_provider_async(self):
        """Get provider and settings for template generation role (async version with landppt support)"""
        logger.info(f"_get_template_role_provider_async called with user_id={self.user_id}")
        
        if self.user_id is not None:
            try:
                from ..db_config_service import get_db_config_service, get_user_ai_provider
                
                config_service = get_db_config_service()
                user_config = await config_service.get_all_config(user_id=self.user_id)
                
                # Get template-specific or default provider and model (use the same mapping as AIConfig)
                role_provider_key, role_model_key = ai_config.MODEL_ROLE_FIELDS.get(
                    "template",
                    ("template_generation_model_provider", "template_generation_model_name"),
                )
                provider_name = self.provider_name or user_config.get(role_provider_key) or user_config.get("default_ai_provider") or "landppt"
                model = user_config.get(role_model_key)
                
                if not model:
                    provider_model_key = f"{provider_name}_model"
                    model = user_config.get(provider_model_key)
                
                logger.info(f"Template provider determined: {provider_name}, model: {model}")
                
                settings = {
                    "role": "template",
                    "provider": provider_name,
                    "model": model
                }
                
                # Use get_user_ai_provider which properly handles landppt system-level config
                provider = await get_user_ai_provider(self.user_id, provider_name)
                logger.info(f"Got template provider successfully: {provider_name}, model: {model}")
                return provider, settings
                
            except Exception as e:
                logger.error(f"Failed to get user template provider async: {e}", exc_info=True)
        else:
            logger.warning("user_id is None, falling back to global config")
        
        # Fall back to global config
        logger.warning("Falling back to global get_role_provider")
        return get_role_provider("template", provider_override=self.provider_name)

    async def _text_completion(self, *, prompt: str, **kwargs):
        provider, settings = await self._get_template_role_provider_async()
        if settings.get("model"):
            kwargs.setdefault("model", settings["model"])
        prompt = SystemPrompts.with_text_cache_prefix(prompt)
        return await provider.text_completion(prompt=prompt, **kwargs)

    async def _chat_completion(self, *, messages: List[AIMessage], **kwargs):
        provider, settings = await self._get_template_role_provider_async()
        if settings.get("model"):
            kwargs.setdefault("model", settings["model"])
        messages = SystemPrompts.normalize_messages_for_cache(messages)
        return await provider.chat_completion(messages=messages, **kwargs)

    async def _stream_text_completion(self, *, prompt: str, **kwargs):
        provider, settings = await self._get_template_role_provider_async()
        if settings.get("model"):
            kwargs.setdefault("model", settings["model"])
        prompt = SystemPrompts.with_text_cache_prefix(prompt)
        if hasattr(provider, 'stream_text_completion'):
            async for chunk in provider.stream_text_completion(prompt=prompt, **kwargs):
                yield chunk
        else:
            response = await provider.text_completion(prompt=prompt, **kwargs)
            yield response.content

    async def _stream_chat_completion(self, *, messages: List[AIMessage], **kwargs):
        provider, settings = await self._get_template_role_provider_async()
        if settings.get("model"):
            kwargs.setdefault("model", settings["model"])
        messages = SystemPrompts.normalize_messages_for_cache(messages)
        if hasattr(provider, 'stream_chat_completion'):
            async for chunk in provider.stream_chat_completion(messages=messages, **kwargs):
                yield chunk
        else:
            response = await provider.chat_completion(messages=messages, **kwargs)
            yield response.content

    async def create_template(self, template_data: Dict[str, Any]) -> Dict[str, Any]:
        """Create a new global master template"""
        try:
            template_settings = ai_config.get_model_config_for_role("template", provider_override=self.provider_name)
            # Validate required fields
            required_fields = ['template_name', 'html_template']
            for field in required_fields:
                if not template_data.get(field):
                    raise ValueError(f"Missing required field: {field}")

            # Check if template name already exists
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                existing = await db_service.get_global_master_template_by_name(
                    template_data['template_name'],
                    user_id=self.user_id,
                )
                if existing:
                    raise ValueError(f"Template name '{template_data['template_name']}' already exists")

            # Generate preview image if not provided
            if not template_data.get('preview_image'):
                template_data['preview_image'] = await self._generate_preview_image(template_data['html_template'])

            # Extract style config if not provided
            if not template_data.get('style_config'):
                template_data['style_config'] = self._extract_style_config(template_data['html_template'])

            # Set default values
            template_data.setdefault('description', '')
            template_data.setdefault('tags', [])
            template_data.setdefault('is_default', False)
            template_data.setdefault('is_active', True)
            template_data.setdefault('created_by', 'system')
            if self.user_id is not None:
                template_data.setdefault('user_id', self.user_id)

            # Create template
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                template = await db_service.create_global_master_template(
                    template_data,
                    user_id=self.user_id,
                )

                return {
                    "id": template.id,
                    "user_id": template.user_id,
                    "template_name": template.template_name,
                    "description": template.description,
                    "preview_image": template.preview_image,
                    "tags": template.tags,
                    "is_default": template.is_default,
                    "is_active": template.is_active,
                    "usage_count": template.usage_count,
                    "created_by": template.created_by,
                    "created_at": template.created_at,
                    "updated_at": template.updated_at
                }

        except IntegrityError as e:
            logger.warning(f"Template create integrity constraint violation: {e}")
            raise ValueError(f"Template name '{template_data.get('template_name')}' already exists")
        except Exception as e:
            logger.error(f"Failed to create global master template: {e}")
            raise

    async def get_all_templates(self, active_only: bool = True) -> List[Dict[str, Any]]:
        """Get all global master templates"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                templates = await db_service.get_all_global_master_templates(
                    active_only,
                    user_id=self.user_id,
                )

                return [
                    {
                        "id": template.id,
                        "user_id": template.user_id,
                        "template_name": template.template_name,
                        "description": template.description,
                        "preview_image": template.preview_image,
                        "tags": template.tags,
                        "is_default": template.is_default,
                        "is_active": template.is_active,
                        "usage_count": template.usage_count,
                        "created_by": template.created_by,
                        "created_at": template.created_at,
                        "updated_at": template.updated_at
                    }
                    for template in templates
                ]

        except Exception as e:
            logger.error(f"Failed to get global master templates: {e}")
            raise

    async def get_all_templates_paginated(
        self,
        active_only: bool = True,
        page: int = 1,
        page_size: int = 6,
        search: Optional[str] = None
    ) -> Dict[str, Any]:
        """Get all global master templates with pagination"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)

                # Calculate offset
                offset = (page - 1) * page_size

                # Get templates with pagination
                templates, total_count = await db_service.get_global_master_templates_paginated(
                    active_only=active_only,
                    offset=offset,
                    limit=page_size,
                    search=search,
                    user_id=self.user_id,
                )

                # Calculate pagination info
                total_pages = (total_count + page_size - 1) // page_size
                has_next = page < total_pages
                has_prev = page > 1

                template_list = [
                    {
                        "id": template.id,
                        "user_id": template.user_id,
                        "template_name": template.template_name,
                        "description": template.description,
                        "preview_image": template.preview_image,
                        "tags": template.tags,
                        "is_default": template.is_default,
                        "is_active": template.is_active,
                        "usage_count": template.usage_count,
                        "created_by": template.created_by,
                        "created_at": template.created_at,
                        "updated_at": template.updated_at
                    }
                    for template in templates
                ]

                return {
                    "templates": template_list,
                    "pagination": {
                        "current_page": page,
                        "page_size": page_size,
                        "total_count": total_count,
                        "total_pages": total_pages,
                        "has_next": has_next,
                        "has_prev": has_prev
                    }
                }
        except Exception as e:
            logger.error(f"Failed to get paginated templates: {e}")
            raise

    async def get_template_by_id(self, template_id: int) -> Optional[Dict[str, Any]]:
        """Get global master template by ID"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                template = await db_service.get_global_master_template_by_id(
                    template_id,
                    user_id=self.user_id,
                )

                if not template:
                    return None

                return {
                    "id": template.id,
                    "user_id": template.user_id,
                    "template_name": template.template_name,
                    "description": template.description,
                    "html_template": template.html_template,
                    "preview_image": template.preview_image,
                    "style_config": template.style_config,
                    "tags": template.tags,
                    "is_default": template.is_default,
                    "is_active": template.is_active,
                    "usage_count": template.usage_count,
                    "created_by": template.created_by,
                    "created_at": template.created_at,
                    "updated_at": template.updated_at
                }

        except Exception as e:
            logger.error(f"Failed to get global master template {template_id}: {e}")
            raise

    async def update_template(self, template_id: int, update_data: Dict[str, Any]) -> bool:
        """Update a global master template"""
        try:
            # Check if template name conflicts (if being updated)
            if 'template_name' in update_data:
                async with AsyncSessionLocal() as session:
                    db_service = DatabaseService(session)
                    existing = await db_service.get_global_master_template_by_name(
                        update_data['template_name'],
                        user_id=self.user_id,
                    )
                    if existing and existing.id != template_id:
                        raise ValueError(f"Template name '{update_data['template_name']}' already exists")

            # Update preview image if HTML template is updated
            if 'html_template' in update_data and 'preview_image' not in update_data:
                update_data['preview_image'] = await self._generate_preview_image(update_data['html_template'])

            # Update style config if HTML template is updated
            if 'html_template' in update_data and 'style_config' not in update_data:
                update_data['style_config'] = self._extract_style_config(update_data['html_template'])

            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                return await db_service.update_global_master_template(
                    template_id,
                    update_data,
                    user_id=self.user_id,
                    allow_system_write=self.allow_system_template_write,
                )

        except IntegrityError as e:
            logger.warning(f"Template update integrity constraint violation: {e}")
            raise ValueError(f"Template name '{update_data.get('template_name')}' already exists")
        except Exception as e:
            logger.error(f"Failed to update global master template {template_id}: {e}")
            raise

    async def delete_template(self, template_id: int) -> bool:
        """Delete a global master template"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)

                # Check if template exists
                template = await db_service.get_global_master_template_by_id(
                    template_id,
                    user_id=self.user_id,
                )
                if not template:
                    logger.warning(f"Template {template_id} not found for deletion")
                    return False

                # Check if it's the default template
                if template.is_default:
                    raise ValueError("Cannot delete the default template")

                logger.info(f"Deleting template {template_id}: {template.template_name}")
                result = await db_service.delete_global_master_template(
                    template_id,
                    user_id=self.user_id,
                    allow_system_write=self.allow_system_template_write,
                )

                if result:
                    logger.info(f"Successfully deleted template {template_id}")
                else:
                    logger.warning(f"Failed to delete template {template_id} - no rows affected")

                return result

        except Exception as e:
            logger.error(f"Failed to delete global master template {template_id}: {e}")
            raise

    async def set_default_template(self, template_id: int) -> bool:
        """Set a template as default"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                return await db_service.set_default_global_master_template(
                    template_id,
                    user_id=self.user_id,
                    allow_system_write=self.allow_system_template_write,
                )

        except Exception as e:
            logger.error(f"Failed to set default template {template_id}: {e}")
            raise

    async def get_default_template(self) -> Optional[Dict[str, Any]]:
        """Get the default template"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                template = await db_service.get_default_global_master_template(user_id=self.user_id)

                if not template:
                    return None

                return {
                    "id": template.id,
                    "user_id": template.user_id,
                    "template_name": template.template_name,
                    "description": template.description,
                    "html_template": template.html_template,
                    "preview_image": template.preview_image,
                    "style_config": template.style_config,
                    "tags": template.tags,
                    "is_default": template.is_default,
                    "is_active": template.is_active,
                    "usage_count": template.usage_count,
                    "created_by": template.created_by,
                    "created_at": template.created_at,
                    "updated_at": template.updated_at
                }

        except Exception as e:
            logger.error(f"Failed to get default template: {e}")
            raise

    @staticmethod
    def _get_template_resource_performance_prompt_text() -> str:
        """统一模板生成阶段的资源可达性与性能约束。"""
        return TemplatePrompts.get_template_resource_performance_prompt_text()

    @staticmethod
    def _get_template_annotation_prompt_text() -> str:
        """固定画布与母版职责分层提示。"""
        return TemplatePrompts.get_template_annotation_prompt_text()

    @staticmethod
    def _get_template_generation_creative_prompt_text() -> str:
        """母版创意目标，强调视觉语言系统而非单页样张。"""
        return TemplatePrompts.get_template_generation_creative_prompt_text()

    @staticmethod
    def _get_template_generation_method_prompt_text() -> str:
        """模板生成的创意思考顺序。"""
        return TemplatePrompts.get_template_generation_method_prompt_text()

    def _get_template_generation_requirements_prompt_text(self) -> str:
        """母版生成技术要求。"""
        return TemplatePrompts.get_template_generation_requirements_prompt_text()

    def _build_template_generation_prompt(self, user_prompt: str, mode_instruction: str = "") -> str:
        """组装模板生成提示词。"""
        return TemplatePrompts.build_template_generation_prompt(user_prompt, mode_instruction=mode_instruction)

    async def generate_template_with_ai(self, prompt: str, template_name: str, description: str = "",
                                      tags: List[str] = None, generation_mode: str = "text_only",
                                      reference_image: dict = None, reference_images: List[dict] = None,
                                      reference_pptx: dict = None,
                                      prompt_is_ready: bool = False):
        """Generate a new template using AI (non-streaming) - does not save to database"""
        import json

        extra_images: list = []  # slide screenshot images for multimodal AI

        if generation_mode == "pptx_extract":
            if not reference_pptx:
                raise ValueError("PPTX提取模式需要上传PPTX文件")

            pptx_context = await self._extract_pptx_template_reference(reference_pptx)
            extracted_summary = pptx_context.get("analysis_summary", "")
            extracted_image = pptx_context.get("reference_image")
            uploaded_images = pptx_context.get("uploaded_images", [])

            # Build image resource hints for the AI
            image_hints = ""
            if uploaded_images:
                image_hints = "\n以下是从PPTX中提取并上传到图床的图片资源，请在HTML中直接引用这些URL：\n"
                for img in uploaded_images:
                    role_label = {'logo': 'Logo', 'background': '背景图', 'icon': '图标',
                                  'decoration': '装饰图', 'content': '内容图'}.get(img.get('role', ''), '图片')
                    image_hints += f"- {role_label}: {img['url']}\n"

            prompt = (
                f"{prompt}\n\n"
                "请基于以下从上传PPTX中提取的完整模板信息生成HTML母版模板。\n"
                "务必忠实还原：颜色主题、字体配置、背景样式（渐变/图片/纯色）、Logo与装饰元素的位置。\n"
                "重点提取视觉风格、版式结构、字体与配色规律，不要照搬原始文案内容。\n\n"
                "如果从多页中推断出稳定的母版元素（如页眉、页脚、页码区域），请在生成结果中保留它们的相对位置和风格。\n"
                f"{image_hints}\n"
                f"{extracted_summary}"
            )

            extra_images = pptx_context.get("extra_reference_images", [])

            if extracted_image:
                reference_image = extracted_image
                generation_mode = "reference_style"
            else:
                generation_mode = "text_only"

        # --- Handle multi-image uploads (reference_style / exact_replica) ------
        # First image → reference_image (primary), images[1:] → extra_images
        if reference_images and len(reference_images) > 0 and generation_mode != "pptx_extract":
            # Always use first uploaded image as primary reference
            reference_image = reference_images[0]
            # Remaining images as extra visual references
            for img in reference_images[1:]:
                if img.get('data'):
                    extra_images.append({
                        "filename": img.get('filename', 'ref_image'),
                        "data": img['data'],
                        "type": img.get('type', 'image/png'),
                        "size": img.get('size', 0),
                    })
            logger.info(f"Multi-image upload: {len(reference_images)} images total, 1 primary + {len(extra_images)} extra")

        # 构建AI提示词
        if generation_mode == "text_only" or not reference_image:
            # 纯文本生成模式
            ai_prompt = prompt if prompt_is_ready else self._build_template_generation_prompt(prompt)
            messages = [{"role": "user", "content": ai_prompt}]
        else:
            # 多模态生成模式
            total_images = 1 + len(extra_images)
            if generation_mode == "reference_style":
                if total_images > 1:
                    mode_instruction = f"""
你收到了 {total_images} 张参考截图，分别展示了PPT模板的不同页面类型（可能包含：封面、目录页、章节过渡页、内容页、结尾页等）。
请**逐张仔细观察**每张截图的配色、字体、布局、装饰元素，综合提炼出统一的视觉风格体系。
按 PPT 母版需求重新组织，确保生成的模板能覆盖所有这些页面类型的设计需求。
"""
                else:
                    mode_instruction = """
请参考上传图片的气质、配色和版式逻辑，但按 PPT 母版需求重新组织，不要机械复刻单一页面。
"""
            else:  # exact_replica
                if total_images > 1:
                    mode_instruction = f"""
你收到了 {total_images} 张参考截图，分别对应PPT模板的不同页面类型（封面、目录、章节过渡、内容、结尾等）。
请**逐张精确还原**每张截图中的视觉设计：配色方案、字体层级、元素位置、背景样式、装饰图形。
生成的HTML母版模板必须包含所有截图中展示的页面类型变体。
"""
                else:
                    mode_instruction = """
请尽量贴近上传图片的风格和版式特征，同时保留可复用的标题、内容和页脚结构。
"""

            ai_prompt = prompt if prompt_is_ready else self._build_template_generation_prompt(
                prompt,
                mode_instruction=mode_instruction,
            )

            # 构建多模态消息
            # 确保图片URL格式正确
            image_data = reference_image['data']
            if not image_data.startswith("data:"):
                # 如果是纯base64数据,添加data URL前缀
                image_data = f"data:{reference_image['type']};base64,{image_data}"

            content_parts = [
                {"type": "text", "text": ai_prompt},
                {
                    "type": "image_url",
                    "image_url": {"url": image_data}
                },
            ]

            # Append all extra reference images (user-uploaded or PPTX screenshots)
            for extra_img in extra_images:
                extra_data = extra_img.get('data', '')
                if extra_data:
                    if not extra_data.startswith("data:"):
                        extra_data = f"data:{extra_img.get('type', 'image/png')};base64,{extra_data}"
                    content_parts.append({
                        "type": "image_url",
                        "image_url": {"url": extra_data}
                    })

            messages = [{"role": "user", "content": content_parts}]

        try:
            # 获取模板生成任务的配置（使用异步版本以支持 landppt 系统级配置）
            provider, template_settings = await self._get_template_role_provider_async()

            if not provider:
                raise ValueError("AI服务未配置或不可用")

            # 转换消息格式
            ai_messages = []
            for msg in messages:
                if isinstance(msg["content"], str):
                    # 纯文本消息
                    ai_messages.append(AIMessage(
                        role=MessageRole.USER,
                        content=[TextContent(text=msg["content"])]
                    ))
                else:
                    # 多模态消息
                    content_parts = []
                    for part in msg["content"]:
                        if part["type"] == "text":
                            content_parts.append(TextContent(text=part["text"]))
                        elif part["type"] == "image_url":
                            # 提取图片URL (已经是完整的data URL格式)
                            image_url = part["image_url"]["url"]
                            if image_url.startswith("data:"):
                                content_parts.append(ImageContent(
                                    image_url={"url": image_url},
                                    content_type=MessageContentType.IMAGE_URL
                                ))
                    ai_messages.append(AIMessage(
                        role=MessageRole.USER,
                        content=content_parts
                    ))

            ai_response = await self._chat_completion(
                messages=ai_messages,
                model=template_settings.get('model')
            )
            full_response = ai_response.content

            if not full_response or not full_response.strip():
                raise ValueError("AI服务返回空响应")

            html_template = self._extract_html_from_response(full_response)
            if not html_template or not html_template.strip():
                raise ValueError("AI响应中未找到有效的HTML模板")

            logger.info(f"Generated HTML template length: {len(html_template)}")

            # 返回结果（不保存到数据库）
            return {
                'html_template': html_template,
                'template_name': template_name,
                'description': description or f"AI生成的模板：{prompt[:100]}",
                'tags': tags or ['AI生成'],
                'llm_response': full_response  # 包含完整的LLM响应
            }

        except Exception as e:
            logger.error(f"Failed to generate template with AI: {e}", exc_info=True)
            raise

    async def generate_template_with_ai_stream(self, prompt: str, template_name: str, description: str = "",
                                             tags: List[str] = None, generation_mode: str = "text_only",
                                             reference_image: dict = None, prompt_is_ready: bool = False):
        """Generate a new template using AI with streaming response"""
        import asyncio
        import json

        # 构建AI提示词
        if generation_mode == "text_only" or not reference_image:
            # 纯文本生成模式
            ai_prompt = prompt if prompt_is_ready else self._build_template_generation_prompt(prompt)
        else:
            # 多模态生成模式
            if generation_mode == "reference_style":
                mode_instruction = """
请参考图片的气质、配色和版式逻辑，但按 PPT 母版需求重新组织，不要机械复刻单一页面。
"""
            else:  # exact_replica
                mode_instruction = """
请尽量贴近参考图片的风格和版式特征，同时保留可复用的标题、内容和页脚结构。
"""

            ai_prompt = prompt if prompt_is_ready else self._build_template_generation_prompt(
                prompt,
                mode_instruction=mode_instruction,
            )

        try:
            # 获取模板生成任务的配置（使用异步版本以支持 landppt 系统级配置）
            provider, template_settings = await self._get_template_role_provider_async()

            # 构建AI消息
            if generation_mode != "text_only" and reference_image:
                # 多模态消息
                # 确保图片URL格式正确 (OpenAI需要完整的data URL格式)
                image_url = reference_image["data"]
                if not image_url.startswith("data:"):
                    # 如果是纯base64数据,添加data URL前缀
                    image_type = reference_image.get("type", "image/png")
                    image_url = f"data:{image_type};base64,{image_url}"

                content_parts = [
                    TextContent(text=ai_prompt),
                    ImageContent(image_url={"url": image_url})
                ]
                messages = [AIMessage(role=MessageRole.USER, content=content_parts)]

                # 检查AI提供商是否支持流式聊天
                if hasattr(provider, 'stream_chat_completion'):
                    # 使用流式聊天API
                    full_response = ""
                    async for chunk in provider.stream_chat_completion(
                        messages=messages,
                        temperature=0.7,
                        model=template_settings.get('model')
                    ):
                        full_response += chunk
                        yield {
                            'type': 'thinking',
                            'content': chunk
                        }
                else:
                    # 使用标准聊天API
                    response = await self._chat_completion(
                        messages=messages,
                        temperature=0.7
                    )
                    full_response = response.content

                    # 模拟流式输出
                    yield {'type': 'thinking', 'content': '🖼️ 正在分析参考图片...\n\n'}
                    await asyncio.sleep(1)
                    yield {'type': 'thinking', 'content': full_response}
            else:
                # 纯文本消息
                if hasattr(provider, 'stream_text_completion'):
                    # 使用流式API
                    full_response = ""
                    async for chunk in provider.stream_text_completion(
                        prompt=ai_prompt,
                        temperature=0.7,
                        model=template_settings.get('model')
                    ):
                        full_response += chunk
                        yield {
                            'type': 'thinking',
                            'content': chunk
                        }
                else:
                    # 使用标准文本完成API
                    response = await provider.text_completion(
                        prompt=ai_prompt,
                        temperature=0.7,
                        model=template_settings.get('model')
                    )
                    full_response = response.content

                    # 模拟流式输出
                    yield {'type': 'thinking', 'content': '🤔 正在分析您的需求...\n\n'}
                    await asyncio.sleep(1)
                    yield {'type': 'thinking', 'content': full_response}

                # 流式完成后，处理完整响应
                yield {'type': 'thinking', 'content': '\n\n✨ 优化样式和交互效果...\n'}
                await asyncio.sleep(0.5)

                # 处理AI响应
                html_template = self._extract_html_from_response(full_response)

                if not html_template or not html_template.strip():
                    raise ValueError("Generated HTML template is empty")

                yield {'type': 'thinking', 'content': '✅ 模板生成完成，准备预览...\n'}
                await asyncio.sleep(0.3)

                # 返回生成完成的信息，包含HTML模板用于预览
                yield {
                    'type': 'complete',
                    'message': '模板生成完成！',
                    'html_template': html_template,
                    'template_name': template_name,
                    'description': description or f"AI生成的模板：{prompt[:100]}",
                    'tags': tags or ['AI生成'],
                    'llm_response': full_response  # 添加完整的LLM响应
                }

        except Exception as e:
            logger.error(f"Failed to generate template with AI stream: {e}", exc_info=True)
            yield {
                'type': 'error',
                'message': str(e)
            }

    async def adjust_template_with_ai_stream(self, current_html: str, adjustment_request: str, template_name: str = "模板"):
        """根据用户反馈调整现有模板，允许做同源升级而不是表面修补。"""
        import asyncio

        ai_prompt = f"""
作为专业的 PPT 模板设计师，请根据用户的调整需求升级现有 HTML 模板。

当前模板：
```html
{current_html}
```

用户调整需求：{adjustment_request}

请按以下要求调整：
- 保留原模板的稳定锚点、占位符、视觉母语和可复用骨架。
- 如果用户需求涉及主舞台表达、内容承载能力、节奏、材质或构图气质，允许对主舞台区进行同源重设计，不要只做表面微调。
- 标题锚点和编号锚点可以做同源优化，但不要改成另一套完全不同的系统；编号锚点不默认在 footer，可位于模板已有的任何稳定位置。
- 如果用户明确要求改变风格方向，可以保留职责分层和占位符，同时重建色彩、字体、材质和组件语气。
- 不要把修改做成局部补丁拼接；调整后的结果仍应像一套完整的新版本母板，而不是在旧页面上缝补。
- 保留所有占位符与可复用结构。
- 输出完整 HTML，继续满足 1280x720、无滚动条，并在固定画布内稳定成立。
{self._get_template_resource_performance_prompt_text()}
{self._get_template_generation_creative_prompt_text()}
{self._get_template_annotation_prompt_text()}

请直接输出调整后的完整 HTML，使用```html```代码块返回。
"""

        try:
            provider, template_settings = await self._get_template_role_provider_async()

            if hasattr(provider, 'stream_text_completion'):
                full_response = ""
                async for chunk in provider.stream_text_completion(
                    prompt=ai_prompt,
                    temperature=0.7,
                    model=template_settings.get('model')
                ):
                    full_response += chunk
                    yield {
                        'type': 'thinking',
                        'content': chunk
                    }

                yield {'type': 'thinking', 'content': '\n\n完成模板调整...\n'}
                await asyncio.sleep(0.5)

                html_template = self._extract_html_from_response(full_response)

                if not html_template or not html_template.strip():
                    raise ValueError("Adjusted HTML template is empty")

                yield {
                    'type': 'complete',
                    'message': '模板调整完成',
                    'html_template': html_template,
                    'template_name': template_name
                }

            else:
                yield {'type': 'thinking', 'content': '正在分析调整需求...\n\n'}
                await asyncio.sleep(1)

                yield {'type': 'thinking', 'content': f'调整需求：{adjustment_request}\n\n'}
                await asyncio.sleep(0.5)

                yield {'type': 'thinking', 'content': '开始重组模板语言...\n'}
                await asyncio.sleep(1)

                response = await provider.text_completion(
                    prompt=ai_prompt,
                    temperature=0.7,
                    model=template_settings.get('model')
                )

                yield {'type': 'thinking', 'content': '完成模板调整...\n'}
                await asyncio.sleep(0.5)

                html_template = self._extract_html_from_response(response.content)

                if not html_template or not html_template.strip():
                    raise ValueError("Adjusted HTML template is empty")

                yield {
                    'type': 'complete',
                    'message': '模板调整完成',
                    'html_template': html_template,
                    'template_name': template_name
                }

        except Exception as e:
            logger.error(f"Failed to adjust template with AI stream: {e}", exc_info=True)
            yield {
                'type': 'error',
                'message': str(e)
            }

    def _decode_uploaded_base64_file(self, raw_data: str) -> bytes:
        """Decode uploaded base64 data (supports raw base64 or data URL)."""
        if not raw_data or not isinstance(raw_data, str):
            raise ValueError("上传文件数据为空")

        data_str = raw_data.strip()
        if data_str.startswith("data:"):
            comma_index = data_str.find(",")
            if comma_index < 0:
                raise ValueError("上传文件数据格式无效")
            data_str = data_str[comma_index + 1 :]

        try:
            return base64.b64decode(data_str, validate=False)
        except Exception as e:
            raise ValueError(f"上传文件Base64解码失败: {e}") from e

    def _safe_pptx_rgb_hex(self, color_obj) -> Optional[str]:
        """Best-effort conversion of python-pptx color object to hex string."""
        try:
            if color_obj is None:
                return None
            rgb = getattr(color_obj, "rgb", None)
            if rgb is None:
                return None
            rgb_text = str(rgb).strip().replace("#", "").replace("0x", "").upper()
            if len(rgb_text) == 6 and all(ch in "0123456789ABCDEF" for ch in rgb_text):
                return f"#{rgb_text}"
        except Exception:
            return None
        return None

    # ------------------------------------------------------------------
    # PPTX deep-extraction helpers
    # ------------------------------------------------------------------

    def _iter_shapes_recursive(self, shapes):
        """Yield all shapes, recursing into group shapes."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        for shape in shapes:
            shape_type = getattr(shape, "shape_type", None)
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                group_shapes = getattr(shape, "shapes", None)
                if group_shapes:
                    yield from self._iter_shapes_recursive(group_shapes)
            else:
                yield shape

    def _extract_theme_colors(self, prs) -> List[Dict[str, str]]:
        """Extract the full theme color palette from the presentation theme XML."""
        colors: List[Dict[str, str]] = []
        try:
            from lxml import etree
            theme = prs.slide_masters[0].slide_layouts[0].slide_master.element
            # Navigate to the theme element
            nsmap = {
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
            }
            # Try to get theme from slide master's relationship
            master = prs.slide_masters[0]
            theme_part = None
            for rel in master.part.rels.values():
                if 'theme' in str(getattr(rel, 'reltype', '')):
                    theme_part = rel.target_part
                    break
            if theme_part is None:
                return colors
            theme_xml = theme_part.element
            # Extract clrScheme
            clr_scheme = theme_xml.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme')
            if clr_scheme is not None:
                color_roles = [
                    'dk1', 'dk2', 'lt1', 'lt2',
                    'accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6',
                    'hlink', 'folHlink',
                ]
                role_labels = {
                    'dk1': '深色1', 'dk2': '深色2', 'lt1': '浅色1', 'lt2': '浅色2',
                    'accent1': '强调色1', 'accent2': '强调色2', 'accent3': '强调色3',
                    'accent4': '强调色4', 'accent5': '强调色5', 'accent6': '强调色6',
                    'hlink': '超链接', 'folHlink': '已访问链接',
                }
                ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                for role in color_roles:
                    el = clr_scheme.find(f'{{{ns}}}{role}')
                    if el is not None:
                        # srgbClr or sysClr
                        srgb = el.find(f'{{{ns}}}srgbClr')
                        if srgb is not None:
                            val = srgb.get('val', '')
                            if val:
                                colors.append({'role': role, 'label': role_labels.get(role, role), 'hex': f'#{val.upper()}'})
                        else:
                            sys_clr = el.find(f'{{{ns}}}sysClr')
                            if sys_clr is not None:
                                last_clr = sys_clr.get('lastClr', '')
                                if last_clr:
                                    colors.append({'role': role, 'label': role_labels.get(role, role), 'hex': f'#{last_clr.upper()}'})
            # Extract font scheme
            font_scheme = theme_xml.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}fontScheme')
            if font_scheme is not None:
                ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                for font_type in ['majorFont', 'minorFont']:
                    font_el = font_scheme.find(f'{{{ns}}}{font_type}')
                    if font_el is not None:
                        latin = font_el.find(f'{{{ns}}}latin')
                        ea = font_el.find(f'{{{ns}}}ea')
                        latin_name = latin.get('typeface', '') if latin is not None else ''
                        ea_name = ea.get('typeface', '') if ea is not None else ''
                        label = '主标题字体' if font_type == 'majorFont' else '正文字体'
                        if latin_name or ea_name:
                            colors.append({'role': font_type, 'label': label, 'hex': '', 'font_latin': latin_name, 'font_ea': ea_name})
        except Exception as e:
            logger.debug(f"Theme color extraction partial failure: {e}")
        return colors

    def _extract_gradient_info(self, fill) -> Optional[Dict[str, Any]]:
        """Extract gradient fill information from a shape/background fill."""
        try:
            fill_type = getattr(fill, 'type', None)
            if fill_type is None:
                return None
            from pptx.enum.dml import MSO_THEME_COLOR
            # Check fill type enum value
            type_name = str(fill_type)
            if 'GRADIENT' not in type_name.upper():
                return None
            gradient_stops = []
            try:
                gs_lst = fill._fill.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gs')
                for gs in gs_lst:
                    pos = gs.get('pos', '')
                    srgb = gs.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    color_hex = None
                    if srgb is not None:
                        color_hex = f"#{srgb.get('val', '').upper()}"
                    else:
                        scheme = gs.find('{http://schemas.openxmlformats.org/drawingml/2006/main}schemeClr')
                        if scheme is not None:
                            color_hex = f"scheme:{scheme.get('val', '')}"
                    pos_pct = f"{int(pos) / 1000:.0f}%" if pos else ''
                    gradient_stops.append({'position': pos_pct, 'color': color_hex or 'unknown'})
            except Exception:
                pass
            if gradient_stops:
                return {'type': 'gradient', 'stops': gradient_stops}
        except Exception:
            pass
        return None

    def _extract_background_info(self, background, slide_width: int, slide_height: int) -> Dict[str, Any]:
        """Extract comprehensive background info: solid color, gradient, or image."""
        result: Dict[str, Any] = {'type': 'none'}
        try:
            fill = background.fill
            fill_type = getattr(fill, 'type', None)
            type_name = str(fill_type) if fill_type is not None else ''

            if 'SOLID' in type_name.upper():
                color_hex = self._safe_pptx_rgb_hex(getattr(fill, 'fore_color', None))
                if color_hex:
                    result = {'type': 'solid', 'color': color_hex}
            elif 'GRADIENT' in type_name.upper():
                grad_info = self._extract_gradient_info(fill)
                if grad_info:
                    result = grad_info
            elif 'PICTURE' in type_name.upper() or 'BACKGROUND' in type_name.upper():
                # Background image fill
                try:
                    bg_element = background._element
                    blip_els = bg_element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                    for blip in blip_els:
                        embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if embed:
                            result = {'type': 'image', 'rId': embed}
                            break
                except Exception:
                    pass
            elif 'PATTERNED' in type_name.upper():
                fg = self._safe_pptx_rgb_hex(getattr(fill, 'fore_color', None))
                bg = self._safe_pptx_rgb_hex(getattr(fill, 'back_color', None))
                result = {'type': 'pattern', 'fore_color': fg, 'back_color': bg}

            # Fallback: try XML-level blipFill detection for image backgrounds
            if result['type'] == 'none':
                try:
                    bg_el = background._element
                    blip_fills = bg_el.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                    for blip in blip_fills:
                        embed = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if embed:
                            result = {'type': 'image', 'rId': embed}
                            break
                except Exception:
                    pass
        except Exception as e:
            logger.debug(f"Background extraction partial failure: {e}")
        return result

    def _resolve_bg_image_blob(self, slide_part, rId: str) -> Optional[Dict[str, Any]]:
        """Resolve an image blob from a relationship ID on a slide part."""
        try:
            rel = slide_part.rels.get(rId)
            if rel is None:
                return None
            target_part = rel.target_part
            blob = getattr(target_part, 'blob', None)
            content_type = getattr(target_part, 'content_type', 'image/png') or 'image/png'
            if blob and isinstance(blob, (bytes, bytearray)):
                return {'blob': bytes(blob), 'content_type': str(content_type)}
        except Exception as e:
            logger.debug(f"Failed to resolve bg image rId={rId}: {e}")
        return None

    async def _upload_pptx_image_to_hosting(self, image_bytes: bytes, content_type: str, filename: str) -> Optional[str]:
        """Upload an extracted PPTX image to the local image hosting and return its absolute URL."""
        try:
            from ..image.image_service import get_image_service
            from ..image.models import (
                ImageUploadRequest, ImageInfo, ImageMetadata,
                ImageSourceType, ImageProvider, ImageFormat,
            )
            from ..url_service import build_image_url
            import hashlib
            import uuid

            # Determine format
            ext = 'png'
            fmt = ImageFormat.PNG
            ct_lower = content_type.lower()
            if 'jpeg' in ct_lower or 'jpg' in ct_lower:
                ext = 'jpg'
                fmt = ImageFormat.JPEG
            elif 'webp' in ct_lower:
                ext = 'webp'
                fmt = ImageFormat.WEBP
            elif 'gif' in ct_lower:
                ext = 'gif'
                fmt = ImageFormat.GIF
            elif 'bmp' in ct_lower:
                ext = 'bmp'
                fmt = ImageFormat.BMP

            if not filename.endswith(f'.{ext}'):
                filename = f"{filename}.{ext}"

            image_service = get_image_service()
            if not image_service.initialized:
                await image_service.initialize()

            upload_req = ImageUploadRequest(
                filename=filename,
                content_type=content_type,
                file_size=len(image_bytes),
                title=f"PPTX提取: {filename}",
                description="从PPTX模板中提取的资源",
                tags=['pptx_extract', 'template_resource'],
                category='local_storage',
            )

            result = await image_service.upload_image(upload_req, image_bytes)
            if result.success and result.image_info:
                url = build_image_url(
                    result.image_info.image_id,
                    width=getattr(result.image_info.metadata, 'width', None),
                    height=getattr(result.image_info.metadata, 'height', None),
                )
                return url
            else:
                logger.warning(f"Image upload failed for {filename}: {getattr(result, 'message', 'unknown')}")
        except Exception as e:
            logger.warning(f"Failed to upload PPTX image to hosting: {e}")
        return None

    def _classify_image_role(self, left_ratio, top_ratio, width_ratio, height_ratio, slide_width, slide_height, area) -> str:
        """Classify an image's role: background, logo, icon, decoration, or content."""
        total_area = slide_width * slide_height if slide_width and slide_height else 1
        area_ratio = area / total_area if total_area else 0

        # Full-slide or near-full → background
        if area_ratio > 0.6 and width_ratio and width_ratio > 0.8 and height_ratio and height_ratio > 0.7:
            return 'background'

        # Small images in corner zones → logo
        if width_ratio and height_ratio and width_ratio <= 0.20 and height_ratio <= 0.15:
            if (top_ratio is not None and top_ratio <= 0.15) or (top_ratio is not None and top_ratio >= 0.85):
                return 'logo'

        # Very small → icon
        if area_ratio < 0.02 and width_ratio and width_ratio <= 0.10:
            return 'icon'

        # Medium decorative
        if area_ratio < 0.08:
            return 'decoration'

        return 'content'

    # ------------------------------------------------------------------
    # PPTX slide screenshot rendering
    # ------------------------------------------------------------------

    @staticmethod
    def _find_libreoffice_path() -> Optional[str]:
        """Find LibreOffice executable on the system."""
        import shutil
        import platform as _plat

        for cmd in ('soffice', 'libreoffice'):
            path = shutil.which(cmd)
            if path:
                return path

        system = _plat.system().lower()
        if system == 'windows':
            candidates = [
                r'C:\Program Files\LibreOffice\program\soffice.exe',
                r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            ]
        elif system == 'darwin':
            candidates = ['/Applications/LibreOffice.app/Contents/MacOS/soffice']
        else:
            candidates = ['/usr/bin/soffice', '/usr/bin/libreoffice',
                          '/snap/bin/libreoffice']
        import os as _os
        for c in candidates:
            if _os.path.isfile(c):
                return c
        return None

    async def _render_pptx_slides_to_images(
        self,
        pptx_bytes: bytes,
        max_slides: int = 8,
        dpi: int = 150,
    ) -> List[Dict[str, Any]]:
        """Render PPTX slides to PNG images for AI visual reference.

        Pipeline:
          1. LibreOffice headless → PDF
          2. PyMuPDF (fitz) → per-page PNG
        Falls back gracefully if tools are unavailable.

        Returns list of {slide_idx: int, blob: bytes, content_type: str}.
        """
        import asyncio
        import tempfile
        import os

        results: List[Dict[str, Any]] = []

        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.info("PyMuPDF (fitz) not installed – slide screenshots disabled")
            return results

        with tempfile.TemporaryDirectory(prefix='landppt_render_') as tmpdir:
            pptx_path = os.path.join(tmpdir, 'input.pptx')
            with open(pptx_path, 'wb') as f:
                f.write(pptx_bytes)

            pdf_path = os.path.join(tmpdir, 'input.pdf')
            pdf_created = False

            # --- Strategy 1: LibreOffice headless → PDF -----------------------
            lo_path = self._find_libreoffice_path()
            if lo_path:
                try:
                    import subprocess
                    loop = asyncio.get_event_loop()
                    proc_result = await loop.run_in_executor(
                        None,
                        lambda: subprocess.run(
                            [lo_path, '--headless', '--norestore',
                             '--convert-to', 'pdf',
                             '--outdir', tmpdir, pptx_path],
                            capture_output=True, timeout=180,
                        ),
                    )
                    if os.path.isfile(pdf_path) and os.path.getsize(pdf_path) > 0:
                        pdf_created = True
                    else:
                        logger.debug(
                            "LibreOffice did not produce PDF. stderr=%s",
                            (proc_result.stderr or b'')[:500],
                        )
                except Exception as e:
                    logger.debug(f"LibreOffice PPTX→PDF failed: {e}")

            # --- Strategy 2: Apryse SDK → PDF (fallback) ---------------------
            if not pdf_created:
                try:
                    from ..pdf_to_pptx_converter import get_pdf_to_pptx_converter
                    converter = get_pdf_to_pptx_converter()
                    if converter.is_available():
                        from apryse_sdk.PDFNetPython import PDFDoc, Convert as ApryseConvert
                        doc = PDFDoc()
                        ApryseConvert.OfficeToPDF(doc, pptx_path, None)
                        doc.Save(pdf_path, 0)
                        doc.Close()
                        if os.path.isfile(pdf_path) and os.path.getsize(pdf_path) > 0:
                            pdf_created = True
                except Exception as e:
                    logger.debug(f"Apryse PPTX→PDF fallback failed: {e}")

            if not pdf_created:
                logger.info("No PDF converter available – slide screenshots skipped")
                return results

            # --- Render PDF pages to PNG via PyMuPDF --------------------------
            try:
                doc = fitz.open(pdf_path)
                page_count = min(doc.page_count, max_slides)
                zoom = dpi / 72.0
                mat = fitz.Matrix(zoom, zoom)
                for i in range(page_count):
                    page = doc.load_page(i)
                    pix = page.get_pixmap(matrix=mat, alpha=False)
                    png_bytes = pix.tobytes("png")
                    results.append({
                        'slide_idx': i + 1,
                        'blob': png_bytes,
                        'content_type': 'image/png',
                        'width': pix.width,
                        'height': pix.height,
                    })
                doc.close()
            except Exception as e:
                logger.warning(f"PDF page rendering failed: {e}")

        return results

    def _is_page_number_like_text(self, text: str) -> bool:
        """Heuristic for page number/footer numbering text in PPTX."""
        import re

        if not text:
            return False
        normalized = str(text).strip()
        if not normalized:
            return False

        normalized = normalized.replace("／", "/").replace("丨", "/").replace("|", "/")
        normalized = re.sub(r"\s+", "", normalized)

        patterns = [
            r"^\d+$",
            r"^\d+/\d+$",
            r"^第?\d+页$",
            r"^第?\d+/\d+页$",
            r"^p\.?\d+$",
            r"^page\d+$",
            r"^\d+-\d+$",
        ]
        lower_text = normalized.lower()
        return any(re.fullmatch(pattern, lower_text) for pattern in patterns)

    def _summarize_common_master_candidates(
        self,
        candidates: List[Dict[str, Any]],
        sampled_slide_count: int,
    ) -> Dict[str, Any]:
        """Summarize stable repeated header/footer/page-number-like elements across sampled slides."""
        if sampled_slide_count <= 1 or not candidates:
            return {"summary_lines": [], "stable_elements": []}

        required_count = max(2, int(sampled_slide_count * 0.6 + 0.999))
        aggregate: Dict[str, Dict[str, Any]] = {}

        for item in candidates:
            signature = str(item.get("signature") or "")
            if not signature:
                continue
            slide_idx = int(item.get("slide_idx") or 0)
            if slide_idx <= 0:
                continue

            slot = aggregate.get(signature)
            if slot is None:
                slot = {
                    "signature": signature,
                    "zone": item.get("zone") or "footer",
                    "kind": item.get("kind") or "shape",
                    "position_desc": item.get("position_desc") or "",
                    "style_hint": item.get("style_hint") or "",
                    "text_examples": [],
                    "slides": set(),
                }
                aggregate[signature] = slot

            if slide_idx in slot["slides"]:
                continue
            slot["slides"].add(slide_idx)

            text_example = (item.get("text_example") or "").strip()
            if text_example and text_example not in slot["text_examples"]:
                slot["text_examples"].append(text_example)

        zone_order = {"header": 0, "footer": 1, "page_number": 2}
        zone_labels = {
            "header": "页眉区域",
            "footer": "页脚区域",
            "page_number": "页码区域",
        }
        kind_labels = {
            "text": "文本框",
            "picture": "图片",
            "shape": "图形",
        }

        stable_elements: List[Dict[str, Any]] = []
        for _, item in aggregate.items():
            count = len(item["slides"])
            if count < required_count:
                continue
            stable_elements.append(
                {
                    "zone": item["zone"],
                    "kind": item["kind"],
                    "position_desc": item["position_desc"],
                    "style_hint": item["style_hint"],
                    "text_examples": item["text_examples"][:3],
                    "slide_hits": count,
                    "sampled_slide_count": sampled_slide_count,
                }
            )

        stable_elements.sort(
            key=lambda x: (
                zone_order.get(str(x.get("zone")), 99),
                -int(x.get("slide_hits") or 0),
                str(x.get("kind") or ""),
                str(x.get("position_desc") or ""),
            )
        )

        if not stable_elements:
            return {"summary_lines": [], "stable_elements": []}

        summary_lines = [
            "- 推断公共母版元素（跨采样页稳定出现，建议在HTML模板中保留相对位置与样式）:"
        ]
        for item in stable_elements[:8]:
            zone_label = zone_labels.get(str(item.get("zone")), str(item.get("zone") or "未知区域"))
            kind_label = kind_labels.get(str(item.get("kind")), str(item.get("kind") or "元素"))
            count = int(item.get("slide_hits") or 0)
            position_desc = str(item.get("position_desc") or "")
            style_hint = str(item.get("style_hint") or "").strip()
            text_examples = item.get("text_examples") or []

            line = f"  - {zone_label}：{kind_label} {position_desc}（{count}/{sampled_slide_count}页）"
            if style_hint:
                line += f"；样式={style_hint}"
            if text_examples and str(item.get('zone')) != "page_number":
                line += f"；示例文本={text_examples[0][:40]}"
            if str(item.get("zone")) == "page_number":
                line += "；疑似页码位置"
            summary_lines.append(line)

        return {"summary_lines": summary_lines, "stable_elements": stable_elements}

    async def _extract_pptx_template_reference(self, reference_pptx: Dict[str, Any]) -> Dict[str, Any]:
        """Extract reusable style/layout hints from an uploaded PPTX for AI template generation.

        Comprehensive extraction covers:
        - Theme color palette (clrScheme) and font scheme (majorFont/minorFont)
        - Slide master / slide layout names
        - Per-slide backgrounds: solid, gradient, pattern, image (uploaded to hosting)
        - All images classified as logo / icon / background / decoration / content
        - Font styles: name, size, bold, italic, underline, color
        - Paragraph alignment
        - Vector shape types and fill colors/gradients
        - Group shape recursion
        - Common master element detection (header/footer/page-number)
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except Exception as e:
            raise ValueError(f"当前环境未安装 python-pptx，无法解析PPTX: {e}") from e

        pptx_bytes = self._decode_uploaded_base64_file(reference_pptx.get("data", ""))
        if not pptx_bytes:
            raise ValueError("PPTX文件为空")

        if len(pptx_bytes) > 50 * 1024 * 1024:
            raise ValueError("PPTX文件过大，请控制在 50MB 以内")

        try:
            prs = Presentation(BytesIO(pptx_bytes))
        except Exception as e:
            raise ValueError(f"PPTX解析失败，请确认文件有效: {e}") from e

        slide_count = len(prs.slides)
        if slide_count == 0:
            raise ValueError("PPTX中没有幻灯片，无法提取模板")

        slide_width = int(getattr(prs, "slide_width", 0) or 0)
        slide_height = int(getattr(prs, "slide_height", 0) or 0)
        emu_per_inch = 914400
        slide_width_px = int(round(slide_width * 96 / emu_per_inch)) if slide_width else 0
        slide_height_px = int(round(slide_height * 96 / emu_per_inch)) if slide_height else 0

        # --- Theme extraction ------------------------------------------------
        theme_colors = self._extract_theme_colors(prs)

        # --- Slide master / layout names -------------------------------------
        layout_names: List[str] = []
        try:
            for master in prs.slide_masters:
                for layout in master.slide_layouts:
                    name = getattr(layout, 'name', None)
                    if name and name not in layout_names:
                        layout_names.append(name)
        except Exception:
            pass

        # --- Decide which slides to sample (diverse selection) ---------------
        all_slides = list(prs.slides)
        max_sample = min(slide_count, 8)
        if slide_count <= max_sample:
            sampled_slides = all_slides
        else:
            # Always include first and last; spread rest evenly
            indices = [0, slide_count - 1]
            step = max(1, (slide_count - 2) // (max_sample - 2))
            for i in range(1, slide_count - 1, step):
                if len(indices) >= max_sample:
                    break
                if i not in indices:
                    indices.append(i)
            indices = sorted(set(indices))[:max_sample]
            sampled_slides = [all_slides[i] for i in indices]

        # --- Counters --------------------------------------------------------
        font_counter: Counter = Counter()
        color_counter: Counter = Counter()
        font_size_counter: Counter = Counter()
        font_style_counter: Counter = Counter()  # bold / italic / underline
        alignment_counter: Counter = Counter()
        layout_counter: Counter = Counter()
        slide_summaries: List[str] = []
        common_master_candidates: List[Dict[str, Any]] = []
        background_summaries: List[str] = []

        # Images to upload: list of {blob, content_type, filename, role, slide_idx, bbox_desc}
        images_to_upload: List[Dict[str, Any]] = []
        best_picture = None
        best_picture_area = 0

        # Dedup image blobs to avoid re-uploading identical resources
        _seen_image_hashes: set = set()
        import hashlib as _hl

        for slide_idx, slide in enumerate(sampled_slides, start=1):
            text_box_count = 0
            picture_count = 0
            vector_shape_count = 0
            slide_title = ""
            shape_lines: List[str] = []
            seen_master_signatures: set = set()

            # --- Background ---------------------------------------------------
            bg_info: Dict[str, Any] = {'type': 'none'}
            try:
                bg_info = self._extract_background_info(slide.background, slide_width, slide_height)
            except Exception:
                pass

            bg_desc = ""
            if bg_info['type'] == 'solid':
                bg_desc = f"纯色({bg_info.get('color', '')})"
                color_counter[bg_info.get('color', '')] += 1
            elif bg_info['type'] == 'gradient':
                stops = bg_info.get('stops', [])
                stop_str = " → ".join(s.get('color', '?') for s in stops[:4])
                bg_desc = f"渐变({stop_str})"
                for s in stops:
                    c = s.get('color', '')
                    if c and not c.startswith('scheme:'):
                        color_counter[c] += 1
            elif bg_info['type'] == 'pattern':
                bg_desc = f"图案(前景={bg_info.get('fore_color')}, 背景={bg_info.get('back_color')})"
            elif bg_info['type'] == 'image':
                rId = bg_info.get('rId', '')
                if rId:
                    blob_info = self._resolve_bg_image_blob(slide.part, rId)
                    if blob_info:
                        img_hash = _hl.md5(blob_info['blob'][:4096]).hexdigest()
                        if img_hash not in _seen_image_hashes:
                            _seen_image_hashes.add(img_hash)
                            images_to_upload.append({
                                'blob': blob_info['blob'],
                                'content_type': blob_info['content_type'],
                                'filename': f"slide{slide_idx}_bg",
                                'role': 'background',
                                'slide_idx': slide_idx,
                                'bbox_desc': '全幅背景',
                            })
                        bg_desc = "背景图片(已提取)"
                    else:
                        bg_desc = "背景图片(无法提取)"

            if bg_desc:
                background_summaries.append(f"第{slide_idx}页背景：{bg_desc}")

            # --- Shapes (recursive into groups) --------------------------------
            all_shapes = list(self._iter_shapes_recursive(slide.shapes))
            shape_count = len(all_shapes)

            for shape in all_shapes:
                try:
                    shape_type = getattr(shape, "shape_type", None)
                    left = int(getattr(shape, "left", 0) or 0)
                    top = int(getattr(shape, "top", 0) or 0)
                    width = int(getattr(shape, "width", 0) or 0)
                    height = int(getattr(shape, "height", 0) or 0)
                    bbox_desc = ""
                    if slide_width and slide_height:
                        bbox_desc = (
                            f"x={left / slide_width:.0%},y={top / slide_height:.0%},"
                            f"w={width / slide_width:.0%},h={height / slide_height:.0%}"
                        )
                    left_ratio = (left / slide_width) if slide_width else None
                    top_ratio = (top / slide_height) if slide_height else None
                    width_ratio = (width / slide_width) if slide_width else None
                    height_ratio = (height / slide_height) if slide_height else None
                    right_ratio = (left_ratio + width_ratio) if left_ratio is not None and width_ratio is not None else None
                    bottom_ratio = (top_ratio + height_ratio) if top_ratio is not None and height_ratio is not None else None

                    # ===== PICTURE shapes =====================================
                    if shape_type == MSO_SHAPE_TYPE.PICTURE:
                        picture_count += 1
                        area = max(0, width) * max(0, height)
                        role = self._classify_image_role(
                            left_ratio, top_ratio, width_ratio, height_ratio,
                            slide_width, slide_height, area,
                        )

                        try:
                            image = shape.image
                            content_type = getattr(image, "content_type", "") or ""
                            blob = getattr(image, "blob", None)
                            if blob and isinstance(blob, (bytes, bytearray)) and content_type.startswith("image/"):
                                img_hash = _hl.md5(blob[:4096]).hexdigest()
                                if img_hash not in _seen_image_hashes:
                                    _seen_image_hashes.add(img_hash)
                                    images_to_upload.append({
                                        'blob': bytes(blob),
                                        'content_type': content_type,
                                        'filename': getattr(image, "filename", None) or f"slide{slide_idx}_{role}",
                                        'role': role,
                                        'slide_idx': slide_idx,
                                        'bbox_desc': bbox_desc,
                                    })
                                if area > best_picture_area:
                                    best_picture_area = area
                                    best_picture = {
                                        "blob": bytes(blob),
                                        "content_type": content_type,
                                        "filename": getattr(image, "filename", None) or f"slide_{slide_idx}_image",
                                        "area": area,
                                    }
                        except Exception:
                            pass

                        role_label = {'background': '背景图', 'logo': 'Logo', 'icon': '图标',
                                      'decoration': '装饰图', 'content': '内容图'}.get(role, '图片')
                        if len(shape_lines) < 12:
                            shape_lines.append(f"- {role_label}：{bbox_desc}")

                        # Master candidate detection for header/footer images
                        if (top_ratio is not None and bottom_ratio is not None
                                and left_ratio is not None and width_ratio is not None
                                and height_ratio is not None):
                            zone = None
                            if top_ratio <= 0.22 and bottom_ratio <= 0.35:
                                zone = "header"
                            elif bottom_ratio >= 0.78 and top_ratio >= 0.58:
                                zone = "footer"
                            if zone:
                                sig = (
                                    f"{zone}|picture|"
                                    f"x{round(left_ratio / 0.04) * 0.04:.2f}|"
                                    f"y{round(top_ratio / 0.04) * 0.04:.2f}|"
                                    f"w{round(width_ratio / 0.04) * 0.04:.2f}|"
                                    f"h{round(height_ratio / 0.04) * 0.04:.2f}"
                                )
                                if sig not in seen_master_signatures:
                                    seen_master_signatures.add(sig)
                                    common_master_candidates.append({
                                        "slide_idx": slide_idx, "zone": zone, "kind": "picture",
                                        "signature": sig, "position_desc": bbox_desc or "top/bottom-region",
                                        "style_hint": f"role={role}", "text_example": "",
                                    })
                        continue

                    # ===== TEXT FRAME shapes ===================================
                    if getattr(shape, "has_text_frame", False):
                        text_box_count += 1
                        text = ""
                        try:
                            text = (shape.text or "").strip()
                        except Exception:
                            text = ""

                        if text and not slide_title:
                            try:
                                is_ph = bool(getattr(shape, "is_placeholder", False))
                            except Exception:
                                is_ph = False
                            if is_ph or top < slide_height * 0.25:
                                slide_title = text.splitlines()[0][:80]

                        text_preview = text.replace("\r", " ").replace("\n", " ").strip()[:80] if text else ""

                        # Font / paragraph detail extraction
                        primary_font_name = None
                        primary_font_size = None
                        primary_font_color = None
                        primary_bold = None
                        primary_italic = None
                        primary_underline = None
                        primary_alignment = None
                        try:
                            for paragraph in list(shape.text_frame.paragraphs)[:4]:
                                # Paragraph alignment
                                p_align = getattr(paragraph, "alignment", None)
                                if p_align is not None:
                                    align_name = str(p_align).split(".")[-1] if p_align else None
                                    if align_name:
                                        alignment_counter[align_name] += 1
                                        if primary_alignment is None:
                                            primary_alignment = align_name

                                for run in list(paragraph.runs)[:6]:
                                    font = getattr(run, "font", None)
                                    if not font:
                                        continue
                                    font_name = (getattr(font, "name", None) or "").strip()
                                    if font_name:
                                        font_counter[font_name] += 1
                                        if primary_font_name is None:
                                            primary_font_name = font_name

                                    font_size = getattr(font, "size", None)
                                    if font_size is not None:
                                        try:
                                            pt_value = round(float(font_size.pt), 1)
                                            font_size_counter[f"{pt_value}pt"] += 1
                                            if primary_font_size is None:
                                                primary_font_size = pt_value
                                        except Exception:
                                            pass

                                    font_color = self._safe_pptx_rgb_hex(getattr(font, "color", None))
                                    if font_color:
                                        color_counter[font_color] += 1
                                        if primary_font_color is None:
                                            primary_font_color = font_color

                                    # Bold / Italic / Underline
                                    is_bold = getattr(font, "bold", None)
                                    is_italic = getattr(font, "italic", None)
                                    is_underline = getattr(font, "underline", None)
                                    if is_bold is True:
                                        font_style_counter["bold"] += 1
                                        if primary_bold is None:
                                            primary_bold = True
                                    if is_italic is True:
                                        font_style_counter["italic"] += 1
                                        if primary_italic is None:
                                            primary_italic = True
                                    if is_underline is True or (is_underline is not None and is_underline is not False):
                                        font_style_counter["underline"] += 1
                                        if primary_underline is None:
                                            primary_underline = True
                        except Exception:
                            pass

                        if len(shape_lines) < 12:
                            line = f"- 文本框：{bbox_desc}"
                            style_tags = []
                            if primary_font_name:
                                style_tags.append(primary_font_name)
                            if primary_font_size:
                                style_tags.append(f"{primary_font_size}pt")
                            if primary_bold:
                                style_tags.append("粗体")
                            if primary_italic:
                                style_tags.append("斜体")
                            if primary_underline:
                                style_tags.append("下划线")
                            if primary_font_color:
                                style_tags.append(primary_font_color)
                            if primary_alignment:
                                style_tags.append(f"对齐={primary_alignment}")
                            if style_tags:
                                line += f" [{', '.join(style_tags)}]"
                            if text_preview:
                                line += f" \"{text_preview[:50]}\""
                            shape_lines.append(line)

                        # Master candidate detection
                        if (top_ratio is not None and bottom_ratio is not None
                                and left_ratio is not None and width_ratio is not None
                                and height_ratio is not None):
                            is_top_zone = top_ratio <= 0.22 and bottom_ratio <= 0.35
                            is_bottom_zone = bottom_ratio >= 0.78 and top_ratio >= 0.58
                            is_page_num = (
                                is_bottom_zone
                                and width_ratio <= 0.28
                                and height_ratio <= 0.14
                                and self._is_page_number_like_text(text_preview or text)
                            )

                            zone = None
                            if is_page_num:
                                zone = "page_number"
                            elif is_top_zone:
                                zone = "header"
                            elif is_bottom_zone:
                                zone = "footer"

                            if zone:
                                style_parts = []
                                if primary_font_name:
                                    style_parts.append(f"font={primary_font_name}")
                                if primary_font_size is not None:
                                    style_parts.append(f"size={primary_font_size}pt")
                                if primary_font_color:
                                    style_parts.append(f"color={primary_font_color}")
                                if primary_bold:
                                    style_parts.append("bold")
                                if primary_alignment:
                                    style_parts.append(f"align={primary_alignment}")
                                style_hint = ", ".join(style_parts)

                                sig = (
                                    f"{zone}|text|"
                                    f"x{round(left_ratio / 0.04) * 0.04:.2f}|"
                                    f"y{round(top_ratio / 0.04) * 0.04:.2f}|"
                                    f"w{round(width_ratio / 0.04) * 0.04:.2f}|"
                                    f"h{round(height_ratio / 0.04) * 0.04:.2f}|"
                                    f"fs{round((primary_font_size or 0) / 2) * 2 if primary_font_size else 0}|"
                                    f"fn{(primary_font_name or '').lower()[:24]}"
                                )
                                if sig not in seen_master_signatures:
                                    seen_master_signatures.add(sig)
                                    common_master_candidates.append({
                                        "slide_idx": slide_idx, "zone": zone, "kind": "text",
                                        "signature": sig, "position_desc": bbox_desc or "top/bottom-region",
                                        "style_hint": style_hint,
                                        "text_example": "" if is_page_num else (text_preview or ""),
                                    })
                        continue

                    # ===== VECTOR / AUTO SHAPES ================================
                    vector_shape_count += 1
                    fill_color = None
                    gradient_info = None
                    try:
                        fill = getattr(shape, "fill", None)
                        if fill:
                            fill_color = self._safe_pptx_rgb_hex(getattr(fill, "fore_color", None))
                            if fill_color:
                                color_counter[fill_color] += 1
                            gradient_info = self._extract_gradient_info(fill)
                    except Exception:
                        pass

                    # Line / outline color
                    outline_color = None
                    try:
                        line_obj = getattr(shape, "line", None)
                        if line_obj:
                            outline_color = self._safe_pptx_rgb_hex(getattr(line_obj.color, "rgb", None) if hasattr(line_obj, 'color') else None)
                            if not outline_color:
                                outline_color = self._safe_pptx_rgb_hex(getattr(line_obj, "color", None))
                    except Exception:
                        pass

                    shape_type_name = str(shape_type).split(".")[-1] if shape_type is not None else "UNKNOWN"
                    # Try to get auto-shape name (e.g., ROUNDED_RECTANGLE)
                    auto_shape_name = ""
                    try:
                        auto_shape_type = getattr(shape, "auto_shape_type", None)
                        if auto_shape_type is not None:
                            auto_shape_name = str(auto_shape_type).split(".")[-1]
                    except Exception:
                        pass

                    if len(shape_lines) < 12:
                        desc = f"- 图形({auto_shape_name or shape_type_name})：{bbox_desc}"
                        extras = []
                        if fill_color:
                            extras.append(f"填充={fill_color}")
                        if gradient_info:
                            stops = gradient_info.get('stops', [])
                            extras.append(f"渐变({len(stops)}色)")
                        if outline_color:
                            extras.append(f"边框={outline_color}")
                        if extras:
                            desc += f" [{', '.join(extras)}]"
                        shape_lines.append(desc)

                    # Master candidate for header/footer shapes
                    if (top_ratio is not None and bottom_ratio is not None
                            and left_ratio is not None and width_ratio is not None
                            and height_ratio is not None):
                        zone = None
                        if top_ratio <= 0.22 and bottom_ratio <= 0.35:
                            zone = "header"
                        elif bottom_ratio >= 0.78 and top_ratio >= 0.58:
                            zone = "footer"

                        if zone:
                            sig = (
                                f"{zone}|shape|"
                                f"x{round(left_ratio / 0.04) * 0.04:.2f}|"
                                f"y{round(top_ratio / 0.04) * 0.04:.2f}|"
                                f"w{round(width_ratio / 0.04) * 0.04:.2f}|"
                                f"h{round(height_ratio / 0.04) * 0.04:.2f}|"
                                f"c{(fill_color or '')[:7]}"
                            )
                            if sig not in seen_master_signatures:
                                seen_master_signatures.add(sig)
                                common_master_candidates.append({
                                    "slide_idx": slide_idx, "zone": zone, "kind": "shape",
                                    "signature": sig,
                                    "position_desc": bbox_desc or "top/bottom-region",
                                    "style_hint": f"fill={fill_color}" if fill_color else "",
                                    "text_example": "",
                                })
                except Exception:
                    continue

            layout_counter[f"text={text_box_count},pic={picture_count},shape={vector_shape_count}"] += 1

            slide_summary = [
                f"第{slide_idx}页：",
                f"元素数={shape_count}（文本框{text_box_count}、图片{picture_count}、图形{vector_shape_count}）",
            ]
            if slide_title:
                slide_summary.append(f"标题={slide_title}")
            if bg_desc:
                slide_summary.append(f"背景={bg_desc}")
            if shape_lines:
                slide_summary.append("布局：\n    " + "\n    ".join(shape_lines[:10]))
            slide_summaries.append(" | ".join(slide_summary))

        # --- Upload extracted images to hosting --------------------------------
        uploaded_images: List[Dict[str, str]] = []  # {role, url, slide_idx, bbox_desc}
        # Prioritise: logos first, then backgrounds, then decorations, limit total
        role_priority = {'logo': 0, 'background': 1, 'icon': 2, 'decoration': 3, 'content': 4}
        images_to_upload.sort(key=lambda x: (role_priority.get(x.get('role', 'content'), 99), x.get('slide_idx', 0)))
        upload_limit = 15
        for img_item in images_to_upload[:upload_limit]:
            try:
                url = await self._upload_pptx_image_to_hosting(
                    img_item['blob'],
                    img_item['content_type'],
                    img_item['filename'],
                )
                if url:
                    uploaded_images.append({
                        'role': img_item.get('role', 'content'),
                        'url': url,
                        'slide_idx': img_item.get('slide_idx', 0),
                        'bbox_desc': img_item.get('bbox_desc', ''),
                    })
            except Exception as e:
                logger.debug(f"Failed to upload pptx extracted image: {e}")

        # --- Render per-slide screenshots --------------------------------------
        slide_screenshots: List[Dict[str, Any]] = []  # {slide_idx, url, b64_data}
        try:
            rendered = await self._render_pptx_slides_to_images(
                pptx_bytes, max_slides=len(sampled_slides), dpi=150,
            )
            for rimg in rendered:
                screenshot_entry: Dict[str, Any] = {
                    'slide_idx': rimg['slide_idx'],
                    'width': rimg.get('width', 0),
                    'height': rimg.get('height', 0),
                }
                # Upload to hosting for URL reference in text prompt
                try:
                    url = await self._upload_pptx_image_to_hosting(
                        rimg['blob'], rimg['content_type'],
                        f"slide{rimg['slide_idx']}_screenshot",
                    )
                    if url:
                        screenshot_entry['url'] = url
                except Exception:
                    pass
                # Also keep base64 for multimodal image messages
                if len(rimg['blob']) <= 5 * 1024 * 1024:
                    b64 = base64.b64encode(rimg['blob']).decode('utf-8')
                    screenshot_entry['b64_data'] = f"data:{rimg['content_type']};base64,{b64}"
                    screenshot_entry['content_type'] = rimg['content_type']
                slide_screenshots.append(screenshot_entry)
        except Exception as e:
            logger.warning(f"Slide screenshot rendering failed: {e}")

        # --- Aggregate statistics ----------------------------------------------
        dominant_fonts = [name for name, _ in font_counter.most_common(5)]
        dominant_colors = [color for color, _ in color_counter.most_common(10)]
        dominant_font_sizes = [size for size, _ in font_size_counter.most_common(6)]
        dominant_layouts = [layout for layout, _ in layout_counter.most_common(3)]
        dominant_font_styles = [f"{style}({cnt}次)" for style, cnt in font_style_counter.most_common(3)]
        dominant_alignments = [f"{a}({cnt}次)" for a, cnt in alignment_counter.most_common(4)]
        common_master_summary = self._summarize_common_master_candidates(
            common_master_candidates,
            len(sampled_slides),
        )

        # --- Build reference image for multimodal AI ---------------------------
        # Prefer the first slide screenshot as reference; fall back to best picture
        reference_image = None
        if slide_screenshots and slide_screenshots[0].get('b64_data'):
            first = slide_screenshots[0]
            reference_image = {
                "filename": f"slide1_screenshot.png",
                "size": len(first.get('b64_data', '')),
                "type": first.get('content_type', 'image/png'),
                "data": first['b64_data'],
            }
        elif best_picture and best_picture.get("blob") and len(best_picture["blob"]) <= 10 * 1024 * 1024:
            img_b64 = base64.b64encode(best_picture["blob"]).decode("utf-8")
            content_type = best_picture.get("content_type") or "image/png"
            reference_image = {
                "filename": best_picture.get("filename") or "pptx_reference_image",
                "size": len(best_picture["blob"]),
                "type": content_type,
                "data": f"data:{content_type};base64,{img_b64}",
            }

        # --- Build additional slide images for multimodal AI -------------------
        # Include remaining slide screenshots as extra reference images
        extra_reference_images: List[Dict[str, Any]] = []
        for ss in slide_screenshots[1:]:  # skip first (already used as reference_image)
            if ss.get('b64_data'):
                extra_reference_images.append({
                    "filename": f"slide{ss['slide_idx']}_screenshot.png",
                    "size": len(ss.get('b64_data', '')),
                    "type": ss.get('content_type', 'image/png'),
                    "data": ss['b64_data'],
                    "slide_idx": ss['slide_idx'],
                })

        # --- Build comprehensive summary text ----------------------------------
        summary_parts: List[str] = [
            "【PPTX模板完整提取结果】",
            f"- 文件名：{reference_pptx.get('filename') or 'uploaded.pptx'}",
            f"- 幻灯片总数：{slide_count}（采样{len(sampled_slides)}页）",
        ]
        if slide_width_px and slide_height_px:
            summary_parts.append(f"- 页面尺寸（像素）：{slide_width_px}x{slide_height_px}")

        # Theme info
        if theme_colors:
            theme_color_items = [c for c in theme_colors if c.get('hex')]
            theme_font_items = [c for c in theme_colors if c.get('font_latin') or c.get('font_ea')]
            if theme_color_items:
                palette_str = ", ".join(f"{c['label']}={c['hex']}" for c in theme_color_items)
                summary_parts.append(f"- 主题色板：{palette_str}")
            if theme_font_items:
                for fi in theme_font_items:
                    parts = []
                    if fi.get('font_latin'):
                        parts.append(f"西文={fi['font_latin']}")
                    if fi.get('font_ea'):
                        parts.append(f"中文={fi['font_ea']}")
                    summary_parts.append(f"- {fi['label']}：{', '.join(parts)}")

        if layout_names:
            summary_parts.append(f"- 母版布局类型：{'、'.join(layout_names[:10])}")
        if dominant_layouts:
            summary_parts.append(f"- 常见元素结构：{'；'.join(dominant_layouts)}")
        if dominant_fonts:
            summary_parts.append(f"- 高频字体：{'、'.join(dominant_fonts)}")
        if dominant_font_sizes:
            summary_parts.append(f"- 高频字号：{'、'.join(dominant_font_sizes)}")
        if dominant_font_styles:
            summary_parts.append(f"- 字体样式：{'、'.join(dominant_font_styles)}")
        if dominant_alignments:
            summary_parts.append(f"- 段落对齐：{'、'.join(dominant_alignments)}")
        if dominant_colors:
            summary_parts.append(f"- 高频使用颜色：{'、'.join(dominant_colors)}")

        # Background summary
        if background_summaries:
            summary_parts.append("- 页面背景：")
            for bs in background_summaries[:8]:
                summary_parts.append(f"  {bs}")

        # Slide screenshots
        screenshots_with_url = [ss for ss in slide_screenshots if ss.get('url')]
        if screenshots_with_url:
            summary_parts.append("- 每页幻灯片截图（请仔细参考每页的视觉设计并精确还原）：")
            for ss in screenshots_with_url:
                summary_parts.append(
                    f"  - 第{ss['slide_idx']}页截图：{ss['url']}"
                )
        if slide_screenshots:
            summary_parts.append(
                f"  （共{len(slide_screenshots)}页截图已附在图片消息中，请逐页参考视觉风格）"
            )

        # Uploaded images list
        if uploaded_images:
            role_labels = {'logo': 'Logo资源', 'background': '背景图', 'icon': '图标',
                           'decoration': '装饰图', 'content': '内容图'}
            summary_parts.append("- 提取的图片资源（可在HTML中直接引用这些URL）：")
            for img in uploaded_images:
                label = role_labels.get(img['role'], '图片')
                summary_parts.append(f"  - [{label}] 第{img['slide_idx']}页 {img['bbox_desc']} → {img['url']}")

        if reference_image:
            summary_parts.append("- 已附上幻灯片截图作为视觉参考（附在图片消息中）")

        summary_parts.extend(common_master_summary.get("summary_lines") or [])

        summary_parts.append("- 采样页详细摘要：")
        for line in slide_summaries[:8]:
            summary_parts.append(f"  {line}")

        # Instruction for AI
        summary_parts.append("")
        summary_parts.append("【重要指引】")
        summary_parts.append("- 你收到的图片消息中包含原始PPTX每页的截图，请**逐页仔细观察**截图中的：")
        summary_parts.append("  配色方案、字体大小层级、元素布局位置、背景样式、装饰图形、Logo位置等。")
        summary_parts.append("- 请在生成的HTML模板中**完整复刻**截图中看到的视觉设计，而非仅依赖文字描述。")
        summary_parts.append("- 如果提取到了Logo/图标/背景图等图片资源URL，请在HTML中使用<img>或CSS background-image直接引用这些URL。")
        summary_parts.append("- 保留页眉/页脚/页码等稳定母版元素的相对位置和视觉风格。")
        summary_parts.append("- 忠实还原原始PPTX的设计语言：配色逻辑、字体层级、空间节奏、装饰元素。")

        return {
            "analysis_summary": "\n".join(summary_parts),
            "reference_image": reference_image,
            "extra_reference_images": extra_reference_images,
            "slide_screenshots": slide_screenshots,
            "slide_count": slide_count,
            "sampled_slide_count": len(sampled_slides),
            "stable_master_elements": common_master_summary.get("stable_elements") or [],
            "theme_colors": theme_colors,
            "uploaded_images": uploaded_images,
            "layout_names": layout_names,
        }

    def _extract_html_from_response(self, response_content: str) -> str:
        """Extract HTML code from AI response with improved extraction"""
        import re

        logger.info(f"Extracting HTML from response. Content length: {len(response_content)}")

        # Try to extract HTML code block (most common format)
        html_match = re.search(r'```html\s*(.*?)\s*```', response_content, re.DOTALL)
        if html_match:
            extracted = html_match.group(1).strip()
            logger.info(f"Extracted HTML from code block. Length: {len(extracted)}")
            return extracted

        # Try to extract any code block that contains DOCTYPE
        code_block_match = re.search(r'```[a-zA-Z]*\s*(<!DOCTYPE html.*?</html>)\s*```', response_content, re.DOTALL | re.IGNORECASE)
        if code_block_match:
            extracted = code_block_match.group(1).strip()
            logger.info(f"Extracted HTML from generic code block. Length: {len(extracted)}")
            return extracted

        # Try to extract DOCTYPE HTML directly
        doctype_match = re.search(r'<!DOCTYPE html.*?</html>', response_content, re.DOTALL | re.IGNORECASE)
        if doctype_match:
            extracted = doctype_match.group(0).strip()
            logger.info(f"Extracted HTML from direct match. Length: {len(extracted)}")
            return extracted

        # If no specific pattern found, check if the content itself is HTML
        content_stripped = response_content.strip()
        if content_stripped.lower().startswith('<!doctype html') and content_stripped.lower().endswith('</html>'):
            logger.info(f"Content appears to be direct HTML. Length: {len(content_stripped)}")
            return content_stripped

        # Return original content as last resort
        logger.warning(f"Could not extract HTML from response, returning original content. Preview: {response_content[:200]}")
        return response_content.strip()

    def _validate_html_template(self, html_content: str) -> bool:
        """Validate HTML template with improved error reporting"""
        try:
            if not html_content or not html_content.strip():
                logger.error("HTML validation failed: Content is empty")
                return False

            html_lower = html_content.lower().strip()

            # Check basic HTML structure with more flexible validation
            if not html_lower.startswith('<!doctype html'):
                logger.error(f"HTML validation failed: Missing or incorrect DOCTYPE. Content starts with: {html_content[:100]}")
                return False

            if '</html>' not in html_lower:
                logger.error("HTML validation failed: Missing closing </html> tag")
                return False

            # Check required elements with better error reporting
            required_elements = {
                '<head>': '<head',
                '<body>': '<body',
                '<title>': '<title'
            }
            missing_elements = []

            for element_name, element_pattern in required_elements.items():
                if element_pattern not in html_lower:
                    missing_elements.append(element_name)

            if missing_elements:
                logger.error(f"HTML validation failed: Missing required elements: {missing_elements}")
                return False

            logger.info("HTML template validation passed successfully")
            return True

        except Exception as e:
            logger.error(f"HTML validation failed with exception: {e}")
            return False

    async def _generate_preview_image(self, html_template: str) -> str:
        """Generate preview image for template (placeholder implementation)"""
        # This is a placeholder implementation
        placeholder_svg = """
        <svg width="320" height="180" xmlns="http://www.w3.org/2000/svg">
            <rect width="320" height="180" fill="#f3f4f6"/>
            <text x="160" y="90" text-anchor="middle" font-family="Arial" font-size="14" fill="#6b7280">
                模板预览
            </text>
        </svg>
        """
        return f"data:image/svg+xml;base64,{base64.b64encode(placeholder_svg.encode()).decode()}"

    def _extract_style_config(self, html_content: str) -> Dict[str, Any]:
        """Extract style configuration from HTML"""
        import re

        style_config = {
            "dimensions": "1280x720",
            "aspect_ratio": "16:9",
            "framework": "HTML + CSS"
        }

        try:
            # Extract color configuration
            color_matches = re.findall(r'(?:background|color)[^:]*:\s*([^;]+)', html_content, re.IGNORECASE)
            if color_matches:
                style_config["colors"] = list(set(color_matches[:10]))  # Limit to 10 colors

            # Extract font configuration
            font_matches = re.findall(r'font-family[^:]*:\s*([^;]+)', html_content, re.IGNORECASE)
            if font_matches:
                style_config["fonts"] = list(set(font_matches[:5]))  # Limit to 5 fonts

            # Check for frameworks
            if 'tailwind' in html_content.lower():
                style_config["framework"] = "Tailwind CSS"
            elif 'bootstrap' in html_content.lower():
                style_config["framework"] = "Bootstrap"

        except Exception as e:
            logger.warning(f"Failed to extract style config: {e}")

        return style_config

    async def get_templates_by_tags(self, tags: List[str], active_only: bool = True) -> List[Dict[str, Any]]:
        """Get global master templates by tags"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                templates = await db_service.get_global_master_templates_by_tags(
                    tags,
                    active_only,
                    user_id=self.user_id,
                )

                return [
                    {
                        "id": template.id,
                        "user_id": template.user_id,
                        "template_name": template.template_name,
                        "description": template.description,
                        "preview_image": template.preview_image,
                        "tags": template.tags,
                        "is_default": template.is_default,
                        "is_active": template.is_active,
                        "usage_count": template.usage_count,
                        "created_by": template.created_by,
                        "created_at": template.created_at,
                        "updated_at": template.updated_at
                    }
                    for template in templates
                ]

        except Exception as e:
            logger.error(f"Failed to get global master templates by tags: {e}")
            raise

    async def get_templates_by_tags_paginated(
        self,
        tags: List[str],
        active_only: bool = True,
        page: int = 1,
        page_size: int = 6,
        search: Optional[str] = None
    ) -> Dict[str, Any]:
        """Get global master templates by tags with pagination"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)

                # Calculate offset
                offset = (page - 1) * page_size

                # Get templates with pagination
                templates, total_count = await db_service.get_global_master_templates_by_tags_paginated(
                    tags=tags,
                    active_only=active_only,
                    offset=offset,
                    limit=page_size,
                    search=search,
                    user_id=self.user_id,
                )

                # Calculate pagination info
                total_pages = (total_count + page_size - 1) // page_size
                has_next = page < total_pages
                has_prev = page > 1

                template_list = [
                    {
                        "id": template.id,
                        "user_id": template.user_id,
                        "template_name": template.template_name,
                        "description": template.description,
                        "preview_image": template.preview_image,
                        "tags": template.tags,
                        "is_default": template.is_default,
                        "is_active": template.is_active,
                        "usage_count": template.usage_count,
                        "created_by": template.created_by,
                        "created_at": template.created_at,
                        "updated_at": template.updated_at
                    }
                    for template in templates
                ]

                return {
                    "templates": template_list,
                    "pagination": {
                        "current_page": page,
                        "page_size": page_size,
                        "total_count": total_count,
                        "total_pages": total_pages,
                        "has_next": has_next,
                        "has_prev": has_prev
                    }
                }
        except Exception as e:
            logger.error(f"Failed to get paginated templates by tags: {e}")
            raise

    async def increment_template_usage(self, template_id: int) -> bool:
        """Increment template usage count"""
        try:
            async with AsyncSessionLocal() as session:
                db_service = DatabaseService(session)
                return await db_service.increment_global_master_template_usage(
                    template_id,
                    user_id=self.user_id,
                )

        except Exception as e:
            logger.error(f"Failed to increment template usage {template_id}: {e}")
            raise
